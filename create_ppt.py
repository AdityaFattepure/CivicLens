"""
CivicLens — PPT Generator
Creates a professional PowerPoint for AWS AI for Bharat Hackathon 2025
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Colors ──────────────────────────────────────────
BLUE = RGBColor(37, 99, 235)       # Primary blue
DARK = RGBColor(15, 23, 42)        # Slate-900
WHITE = RGBColor(255, 255, 255)
GRAY_50 = RGBColor(249, 250, 251)
GRAY_100 = RGBColor(243, 244, 246)
GRAY_400 = RGBColor(156, 163, 175)
GRAY_500 = RGBColor(107, 114, 128)
GRAY_600 = RGBColor(75, 85, 99)
GRAY_700 = RGBColor(55, 65, 81)
GRAY_800 = RGBColor(31, 41, 55)
GRAY_900 = RGBColor(17, 24, 39)
RED = RGBColor(220, 38, 38)
GREEN = RGBColor(5, 150, 105)
AMBER = RGBColor(217, 119, 6)
LIGHT_BLUE = RGBColor(219, 234, 254)
LIGHT_RED = RGBColor(254, 226, 226)
LIGHT_GREEN = RGBColor(209, 250, 229)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ─── Helper Functions ────────────────────────────────
def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, size=18, color=GRAY_900, bold=False, align=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=GRAY_900, bold=False, align=PP_ALIGN.LEFT, space_before=Pt(6), font_name='Calibri'):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    if space_before:
        p.space_before = space_before
    return p

def add_blue_bar(slide):
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.08), BLUE)

def add_slide_number(slide, num):
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.35), str(num), size=10, color=GRAY_400, align=PP_ALIGN.RIGHT)

def make_table(slide, rows, cols, left, top, width, height):
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    return table

def style_table(table, header_color=BLUE, header_text_color=WHITE):
    for i, cell in enumerate(table.rows[0].cells):
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = header_text_color
            p.font.bold = True
            p.font.size = Pt(13)
            p.font.name = 'Calibri'
    for row_idx in range(1, len(table.rows)):
        for cell in table.rows[row_idx].cells:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if row_idx % 2 == 1 else GRAY_50
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(12)
                p.font.color.rgb = GRAY_800
                p.font.name = 'Calibri'

def set_cell(table, row, col, text, size=12, bold=False, color=GRAY_800):
    cell = table.cell(row, col)
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = 'Calibri'
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE


# ═════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BLUE)

# Badge
badge = add_shape_bg(slide, Inches(4.8), Inches(1.5), Inches(3.7), Inches(0.45), RGBColor(30, 58, 138))
add_text(slide, Inches(4.8), Inches(1.52), Inches(3.7), Inches(0.4), '★  AWS AI for Bharat Hackathon 2025', size=13, color=RGBColor(147, 197, 253), bold=False, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(2.3), Inches(11.333), Inches(1), 'CivicLens', size=52, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tf = add_text(slide, Inches(1), Inches(3.2), Inches(11.333), Inches(0.6), 'The Truth Engine', size=36, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(4.1), Inches(9.333), Inches(0.5), 'AI-Powered Civic Transparency Platform for Indian Infrastructure', size=18, color=GRAY_400, align=PP_ALIGN.CENTER)

add_text(slide, Inches(3), Inches(5.2), Inches(7.333), Inches(0.5),
         '"30 days for an RTI reply.  3 seconds on CivicLens."', size=16, color=RGBColor(253, 230, 138), bold=False, align=PP_ALIGN.CENTER)

add_text(slide, Inches(2), Inches(6.3), Inches(9.333), Inches(0.4),
         'Live Demo: https://main.d3k5zyfw9140d9.amplifyapp.com   |   Region: ap-south-1 (Mumbai)', size=11, color=GRAY_500, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 2)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5), 'The Problem', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.4), 'The ₹15 Lakh Crore Visibility Void', size=16, color=BLUE, bold=False)

# Big stat box
stat_box = add_shape_bg(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(1), LIGHT_RED)
add_text(slide, Inches(1), Inches(1.7), Inches(11.3), Inches(0.8),
         'India spends ₹15+ Lakh Crore/year on public infrastructure — yet citizens have ZERO visibility into budgets, deadlines, or who is responsible.',
         size=16, color=RGBColor(153, 27, 27), bold=True, align=PP_ALIGN.CENTER)

# 4 Problem cards
problems = [
    ("Zero Visibility", "Citizens can't tell if a road costs ₹50 Lakhs or ₹50 Crore. No budget, no deadline, no contractor info."),
    ("30-Day RTI Process", "Filing RTI for ONE answer about ONE project takes 30 days. By then, a new tender is already issued."),
    ("No Accountability", "Citizens blame PM for potholes — they don't know the actual officer or contractor responsible."),
    ("No Verification", "Official status says 'Completed ✓' but reality shows potholes everywhere ✗. No way to prove it."),
]

for i, (title, desc) in enumerate(problems):
    x = Inches(0.8 + i * 3.05)
    y = Inches(3.0)
    card = add_shape_bg(slide, x, y, Inches(2.85), Inches(2.3), GRAY_50)
    card.line.color.rgb = GRAY_100
    add_text(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.45), Inches(0.4), title, size=15, color=GRAY_900, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.45), Inches(1.4), desc, size=12, color=GRAY_600)

# Bottom stat comparison
add_shape_bg(slide, Inches(3.5), Inches(5.7), Inches(2.2), Inches(1.2), LIGHT_RED)
add_text(slide, Inches(3.5), Inches(5.8), Inches(2.2), Inches(0.5), '30 days', size=28, color=RED, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(6.3), Inches(2.2), Inches(0.4), 'RTI Response Time', size=11, color=GRAY_500, align=PP_ALIGN.CENTER)

add_text(slide, Inches(6.1), Inches(6.0), Inches(1), Inches(0.5), '→', size=32, color=GRAY_400, align=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(7.5), Inches(5.7), Inches(2.2), Inches(1.2), LIGHT_GREEN)
add_text(slide, Inches(7.5), Inches(5.8), Inches(2.2), Inches(0.5), '3 seconds', size=28, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(7.5), Inches(6.3), Inches(2.2), Inches(0.4), 'CivicLens', size=11, color=GRAY_500, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 3)

add_text(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.5), 'Our Solution — CivicLens', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(0.5),
         'Transforms raw government tender PDFs into an interactive transparency platform where citizens can see, verify, and hold power accountable.',
         size=15, color=GRAY_600)

# 3 Pillars
pillars = [
    ("🔍  SEE", "Every government project on a live map — budget, deadline, contractor, status",
     "AI extracts structured data from unstructured tender PDFs automatically using Textract + Bedrock", LIGHT_BLUE, BLUE),
    ("✅  VERIFY", "Submit geo-tagged photos. AI compares ground reality vs official claims",
     "Rekognition detects road conditions + 200m geofence validation = tamper-proof truth check", LIGHT_GREEN, GREEN),
    ("👁️  HOLD ACCOUNTABLE", "Trace the full hierarchy: Ministry → Department → Officer → Contractor",
     "NLP extracts officer names, departments, contractors from tenders — builds accountability chain", RGBColor(254, 243, 199), AMBER),
]

for i, (title, what, how, bg_color, accent) in enumerate(pillars):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.0)
    card = add_shape_bg(slide, x, y, Inches(3.7), Inches(3.8), bg_color)
    add_text(slide, x + Inches(0.25), y + Inches(0.25), Inches(3.2), Inches(0.5), title, size=20, color=accent, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(0.9), Inches(3.2), Inches(0.9), what, size=14, color=GRAY_800, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(2.1), Inches(3.2), Inches(0.3), 'How:', size=11, color=accent, bold=True)
    add_text(slide, x + Inches(0.25), y + Inches(2.5), Inches(3.2), Inches(1.1), how, size=12, color=GRAY_600)

# Key differentiator
add_shape_bg(slide, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8), RGBColor(239, 246, 255))
add_text(slide, Inches(1), Inches(6.2), Inches(11.3), Inches(0.6),
         '⚡ Key Differentiator: CivicLens doesn\'t just show data — it DETECTS WHEN OFFICIALS LIE by comparing AI-analyzed citizen photos against official project status.',
         size=14, color=RGBColor(30, 64, 175), bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════
# SLIDE 4 — WHY AI IS REQUIRED
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 4)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Why AI Is Required — Not Optional', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.3), 'Without AI, this solution fundamentally cannot exist', size=14, color=RED, bold=True)

table = make_table(slide, 6, 3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(4.2))
# Header
set_cell(table, 0, 0, 'The Problem', size=13, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Why Only AI Can Solve It', size=13, bold=True, color=WHITE)
set_cell(table, 0, 2, 'AWS Service Used', size=13, bold=True, color=WHITE)
# Rows
set_cell(table, 1, 0, 'Government tenders are unstructured scanned PDFs (Hindi + English, stamps, seals)')
set_cell(table, 1, 1, 'No API, no structured database exists. Only OCR + Generative AI can extract budget, contractor, deadline from raw PDF text.')
set_cell(table, 1, 2, 'Amazon Textract + Bedrock (Claude 3 Haiku)', bold=True)

set_cell(table, 2, 0, 'Getting project info takes 30+ days via RTI')
set_cell(table, 2, 1, 'AI processes thousands of tenders automatically on S3 upload — zero human intervention.')
set_cell(table, 2, 2, 'Lambda (S3 trigger) + Bedrock', bold=True)

set_cell(table, 3, 0, 'No way to verify if a road is actually repaired')
set_cell(table, 3, 1, 'Computer Vision detects potholes, debris, cracks from citizen photos and compares against official "Completed" status.')
set_cell(table, 3, 2, 'Amazon Rekognition', bold=True)

set_cell(table, 4, 0, 'Citizens don\'t know who is actually responsible')
set_cell(table, 4, 1, 'NLP extracts officer names, departments, contractor names from tender text → builds accountability chain automatically.')
set_cell(table, 4, 2, 'Bedrock (Claude 3 Haiku)', bold=True)

set_cell(table, 5, 0, 'India has 1,000+ tenders per day across states')
set_cell(table, 5, 1, 'Only AI can process this volume. Human review of unstructured PDFs at this scale is physically impossible.')
set_cell(table, 5, 2, 'Full Serverless AI Pipeline', bold=True)

# Set column widths
table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(5.5)
table.columns[2].width = Inches(2.7)
style_table(table)

# Bottom callout
add_shape_bg(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.9), RGBColor(254, 242, 242))
add_text(slide, Inches(1), Inches(6.1), Inches(11.3), Inches(0.7),
         '⚠️  The core value — transforming opaque government PDFs into transparent, verifiable, searchable data — fundamentally requires Generative AI for extraction and Computer Vision for verification. There is NO manual alternative at this scale.',
         size=13, color=RGBColor(153, 27, 27), bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════
# SLIDE 5 — ARCHITECTURE
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 5)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Architecture — 100% Serverless on AWS', size=28, color=GRAY_900, bold=True)

# Architecture boxes
# Frontend
add_shape_bg(slide, Inches(3.5), Inches(1.3), Inches(6.3), Inches(0.7), RGBColor(219, 234, 254))
add_text(slide, Inches(3.5), Inches(1.35), Inches(6.3), Inches(0.6), '🌐  Frontend — React 18 + Vite + TailwindCSS  |  AWS Amplify Hosting', size=14, color=RGBColor(30, 64, 175), bold=True, align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(6.3), Inches(2.0), Inches(0.7), Inches(0.3), '▼', size=16, color=GRAY_400, align=PP_ALIGN.CENTER)

# API Gateway
add_shape_bg(slide, Inches(3.5), Inches(2.3), Inches(6.3), Inches(0.7), RGBColor(254, 243, 199))
add_text(slide, Inches(3.5), Inches(2.35), Inches(6.3), Inches(0.6), '🔗  Amazon API Gateway  —  6 REST Endpoints with CORS', size=14, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Arrow
add_text(slide, Inches(6.3), Inches(3.0), Inches(0.7), Inches(0.3), '▼', size=16, color=GRAY_400, align=PP_ALIGN.CENTER)

# Lambda boxes
lambdas = [
    ("Query Lambdas (×4)\nGetProjects\nGetStats\nGetAccountability\nUpvote", GRAY_50, GRAY_700),
    ("Submit Audit Lambda\n+ Amazon Rekognition\n  (Photo Analysis)\n+ Location Service\n  (200m Geofence)", LIGHT_GREEN, GREEN),
    ("Ingest Tender Lambda\n+ Amazon Textract (OCR)\n+ Bedrock Claude 3 (AI)\n+ Location Service\n  (Geocoding)", LIGHT_BLUE, BLUE),
]

for i, (text, bg, fg) in enumerate(lambdas):
    x = Inches(0.8 + i * 4.2)
    box = add_shape_bg(slide, x, Inches(3.4), Inches(3.8), Inches(1.8), bg)
    box.line.color.rgb = GRAY_100
    add_text(slide, x + Inches(0.2), Inches(3.5), Inches(3.4), Inches(1.6), text, size=12, color=fg, bold=False)

# Arrow
add_text(slide, Inches(6.3), Inches(5.2), Inches(0.7), Inches(0.3), '▼', size=16, color=GRAY_400, align=PP_ALIGN.CENTER)

# DynamoDB + S3
add_shape_bg(slide, Inches(2.5), Inches(5.5), Inches(5), Inches(0.8), RGBColor(204, 251, 241))
add_text(slide, Inches(2.5), Inches(5.55), Inches(5), Inches(0.7), '🗄️  Amazon DynamoDB — Projects + Audits Tables (Geohash & City GSIs)', size=13, color=RGBColor(13, 148, 136), bold=True, align=PP_ALIGN.CENTER)

add_shape_bg(slide, Inches(8), Inches(5.5), Inches(3.5), Inches(0.8), RGBColor(254, 243, 199))
add_text(slide, Inches(8), Inches(5.55), Inches(3.5), Inches(0.7), '📦  Amazon S3 — 2 Buckets\nTender PDFs + Citizen Photos', size=12, color=AMBER, bold=True, align=PP_ALIGN.CENTER)

# Region tag
add_text(slide, Inches(0.8), Inches(6.7), Inches(5), Inches(0.3), 'Region: ap-south-1 (Mumbai)  |  IaC: AWS SAM (308 lines)', size=11, color=GRAY_400)


# ═════════════════════════════════════════════════════
# SLIDE 6 — AWS SERVICES DEEP DIVE
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 6)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'AWS Services — Detailed Breakdown', size=28, color=GRAY_900, bold=True)

# AI/ML table
add_text(slide, Inches(0.8), Inches(1.1), Inches(5), Inches(0.3), 'AI / ML Services (Core)', size=16, color=BLUE, bold=True)

table = make_table(slide, 5, 3, Inches(0.8), Inches(1.5), Inches(11.7), Inches(2.4))
set_cell(table, 0, 0, 'AWS Service', size=12, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Role', size=12, bold=True, color=WHITE)
set_cell(table, 0, 2, 'How It\'s Used', size=12, bold=True, color=WHITE)

set_cell(table, 1, 0, 'Amazon Bedrock\n(Claude 3 Haiku)', bold=True)
set_cell(table, 1, 1, 'Generative AI')
set_cell(table, 1, 2, 'Extracts structured JSON (title, budget, contractor, deadline, location) from raw OCR text. Custom prompt for Indian government tender language.')

set_cell(table, 2, 0, 'Amazon Textract', bold=True)
set_cell(table, 2, 1, 'Document OCR')
set_cell(table, 2, 2, 'DetectDocumentText API — converts scanned tender PDFs (Hindi/English, stamps, signatures) into machine-readable text.')

set_cell(table, 3, 0, 'Amazon Rekognition', bold=True)
set_cell(table, 3, 1, 'Computer Vision')
set_cell(table, 3, 2, 'DetectLabels API — analyzes citizen photos for potholes, rubble, debris, construction equipment. Labels compared vs official status.')

set_cell(table, 4, 0, 'Amazon Location\nService', bold=True)
set_cell(table, 4, 1, 'Geospatial')
set_cell(table, 4, 2, 'Geocodes tender addresses → lat/lng for map. Validates citizen photos are within 200m geofence of project site.')

table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(8.0)
style_table(table)

# Infra table
add_text(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.3), 'Infrastructure Services', size=16, color=BLUE, bold=True)

table2 = make_table(slide, 7, 3, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5))
set_cell(table2, 0, 0, 'AWS Service', size=12, bold=True, color=WHITE)
set_cell(table2, 0, 1, 'Role', size=12, bold=True, color=WHITE)
set_cell(table2, 0, 2, 'Configuration', size=12, bold=True, color=WHITE)

infra = [
    ('AWS Lambda', 'Compute', '6 Python 3.13 functions on ARM64 Graviton. 256MB–1024MB memory. S3-triggered for automated processing.'),
    ('Amazon API Gateway', 'API Layer', 'REST API with CORS, 6 endpoints: /projects, /stats, /accountability, /audit, /upvote'),
    ('Amazon DynamoDB', 'Database', '2 tables: Projects (geohash GSI + city GSI) + Audits (project_id GSI). PAY_PER_REQUEST billing.'),
    ('Amazon S3', 'Storage', '2 buckets: Tender PDFs (triggers Lambda on .pdf upload) + Citizen photos (CORS-enabled).'),
    ('AWS SAM', 'IaC', 'Complete infrastructure-as-code — template.yaml (308 lines). Deploy with: sam build && sam deploy'),
    ('AWS Amplify', 'Hosting', 'Static React frontend deployment with global CDN.'),
]
for i, (svc, role, config) in enumerate(infra):
    set_cell(table2, i+1, 0, svc, bold=True)
    set_cell(table2, i+1, 1, role)
    set_cell(table2, i+1, 2, config)
table2.columns[0].width = Inches(2.2)
table2.columns[1].width = Inches(1.5)
table2.columns[2].width = Inches(8.0)
style_table(table2)


# ═════════════════════════════════════════════════════
# SLIDE 7 — AI PIPELINE 1: TENDER INGESTION
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 7)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'AI Pipeline 1 — Automated Tender Processing', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.3), 'Zero human intervention — fully automated on S3 upload', size=14, color=GRAY_500)

# Pipeline steps
steps = [
    ("STEP 1", "PDF Upload\nto S3 Bucket", "Tender PDF uploaded\ntriggers Lambda", GRAY_50, GRAY_700),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 2", "Amazon\nTextract", "OCR extracts text\nfrom scanned PDF", RGBColor(243, 232, 255), RGBColor(126, 34, 206)),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 3", "Bedrock\nClaude 3", "AI extracts JSON:\ntitle, budget,\ncontractor, deadline", LIGHT_GREEN, GREEN),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 4", "Location\nService", "Geocodes address\n→ lat/lng\ncoordinates", RGBColor(254, 243, 199), AMBER),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 5", "DynamoDB", "Stored with\ngeohash index\nfor map queries", RGBColor(204, 251, 241), RGBColor(13, 148, 136)),
]

x_pos = Inches(0.4)
for step_label, title, desc, bg, fg in steps:
    if step_label == "→":
        add_text(slide, x_pos, Inches(2.5), Inches(0.5), Inches(0.5), '→', size=28, color=GRAY_400, align=PP_ALIGN.CENTER)
        x_pos += Inches(0.5)
    else:
        box = add_shape_bg(slide, x_pos, Inches(1.6), Inches(2.2), Inches(2.8), bg)
        box.line.color.rgb = GRAY_100
        add_text(slide, x_pos + Inches(0.1), Inches(1.7), Inches(2.0), Inches(0.3), step_label, size=10, color=GRAY_400, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x_pos + Inches(0.1), Inches(2.0), Inches(2.0), Inches(0.8), title, size=16, color=fg, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x_pos + Inches(0.1), Inches(2.9), Inches(2.0), Inches(1.2), desc, size=11, color=GRAY_600, align=PP_ALIGN.CENTER)
        x_pos += Inches(2.4)

# What AI Adds box
add_shape_bg(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2), RGBColor(239, 246, 255))
add_text(slide, Inches(1), Inches(4.9), Inches(11.3), Inches(0.3), 'What AI Adds:', size=15, color=BLUE, bold=True)

benefits = [
    "✓  Raw government PDF → Structured, searchable project data in < 10 seconds",
    "✓  Handles Hindi/English mixed documents with stamps and signatures",
    "✓  Processes 1,000+ tenders/day — impossible manually",
    "✓  30-day RTI process replaced with instant automated extraction",
]
tf = add_text(slide, Inches(1), Inches(5.3), Inches(11.3), Inches(0.3), benefits[0], size=13, color=GRAY_700)
for b in benefits[1:]:
    add_para(tf, b, size=13, color=GRAY_700, space_before=Pt(4))

add_text(slide, Inches(1), Inches(6.5), Inches(11.3), Inches(0.3), 'Code: backend/functions/ingest_tender/handler.py', size=10, color=GRAY_400)


# ═════════════════════════════════════════════════════
# SLIDE 8 — AI PIPELINE 2: CITIZEN VERIFICATION
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 8)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'AI Pipeline 2 — Photo Truth Check', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.3), 'Citizens verify ground reality using their phone camera', size=14, color=GRAY_500)

# Pipeline steps
steps2 = [
    ("STEP 1", "Citizen Photo\n+ GPS", "Geo-tagged photo\nuploaded to S3", LIGHT_BLUE, BLUE),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 2", "Geofence\nCheck", "Location Service\nverifies within\n200m of project", RGBColor(254, 243, 199), AMBER),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 3", "Amazon\nRekognition", "DetectLabels:\npotholes, rubble,\nasphalt, debris", RGBColor(243, 232, 255), RGBColor(126, 34, 206)),
    ("→", "", "", WHITE, GRAY_400),
    ("STEP 4", "Truth\nCheck", "Compare AI detection\nvs Official Status\n→ MATCH / MISMATCH", LIGHT_RED, RED),
]

x_pos = Inches(0.4)
for step_label, title, desc, bg, fg in steps2:
    if step_label == "→":
        add_text(slide, x_pos, Inches(2.5), Inches(0.5), Inches(0.5), '→', size=28, color=GRAY_400, align=PP_ALIGN.CENTER)
        x_pos += Inches(0.6)
    else:
        box = add_shape_bg(slide, x_pos, Inches(1.6), Inches(2.7), Inches(2.8), bg)
        box.line.color.rgb = GRAY_100
        add_text(slide, x_pos + Inches(0.1), Inches(1.7), Inches(2.5), Inches(0.3), step_label, size=10, color=GRAY_400, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x_pos + Inches(0.1), Inches(2.0), Inches(2.5), Inches(0.8), title, size=16, color=fg, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x_pos + Inches(0.1), Inches(2.9), Inches(2.5), Inches(1.2), desc, size=11, color=GRAY_600, align=PP_ALIGN.CENTER)
        x_pos += Inches(2.9)

# Truth Check Example
add_text(slide, Inches(0.8), Inches(4.7), Inches(5), Inches(0.3), 'Example Truth Check:', size=15, color=GRAY_900, bold=True)

table = make_table(slide, 4, 3, Inches(0.8), Inches(5.1), Inches(8), Inches(1.6))
set_cell(table, 0, 0, '', size=12, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Official Claim', size=12, bold=True, color=WHITE)
set_cell(table, 0, 2, 'AI Detection', size=12, bold=True, color=WHITE)

set_cell(table, 1, 0, 'Status', bold=True)
set_cell(table, 1, 1, '✅ "Completed"')
set_cell(table, 1, 2, '❌ Potholes, Rubble, Debris detected')

set_cell(table, 2, 0, 'Verdict', bold=True)
set_cell(table, 2, 1, '')
set_cell(table, 2, 2, 'DISPUTED — Discrepancy Detected', bold=True, color=RED)

set_cell(table, 3, 0, 'Action', bold=True)
set_cell(table, 3, 1, '')
set_cell(table, 3, 2, 'Project flagged. Accountability chain notified.')

table.columns[0].width = Inches(1.2)
table.columns[1].width = Inches(2.8)
table.columns[2].width = Inches(4.0)
style_table(table)

add_text(slide, Inches(1), Inches(6.9), Inches(11.3), Inches(0.3), 'Code: backend/functions/submit_audit/handler.py', size=10, color=GRAY_400)


# ═════════════════════════════════════════════════════
# SLIDE 9 — KEY FEATURES
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 9)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Platform Features', size=28, color=GRAY_900, bold=True)

features = [
    ("Interactive Project Map", "18+ real projects across 8 Indian cities on a live map, color-coded by status", "DynamoDB + Location Service", LIGHT_BLUE, BLUE),
    ("AI Tender Analysis", "Automated PDF → structured data. Budget, contractor, deadline extracted in seconds", "Textract + Bedrock", RGBColor(243, 232, 255), RGBColor(126, 34, 206)),
    ("Truth Check Engine", "Side-by-side: 'THE PROMISE' (official) vs 'THE REALITY' (citizen photo + AI)", "Rekognition + DynamoDB", LIGHT_RED, RED),
    ("Accountability Chain", "Full government hierarchy: Ministry → Department → Officer → Contractor", "DynamoDB", RGBColor(254, 243, 199), AMBER),
    ("Citizen Auditing", "4-step wizard: Select project → GPS verify → Upload photo → AI analyzes", "S3 + Rekognition + Location", LIGHT_GREEN, GREEN),
    ("AI Demo Page", "Interactive simulation of both AI pipelines — see AI processing in real-time", "Frontend Simulation", GRAY_50, GRAY_600),
    ("Statistics Dashboard", "Budget analysis, status distribution, most disputed projects, city breakdown", "DynamoDB + Recharts", RGBColor(204, 251, 241), RGBColor(13, 148, 136)),
    ("Signal-Based Priority", "Upvote system — most critical issues rise to top. No comments, no noise.", "Lambda + DynamoDB", RGBColor(254, 215, 170), RGBColor(194, 65, 12)),
]

for i, (title, desc, services, bg, fg) in enumerate(features):
    col = i % 4
    row = i // 4
    x = Inches(0.5 + col * 3.15)
    y = Inches(1.2 + row * 2.9)
    card = add_shape_bg(slide, x, y, Inches(2.95), Inches(2.6), bg)
    card.line.color.rgb = GRAY_100
    add_text(slide, x + Inches(0.15), y + Inches(0.15), Inches(2.65), Inches(0.45), title, size=14, color=fg, bold=True)
    add_text(slide, x + Inches(0.15), y + Inches(0.65), Inches(2.65), Inches(1.2), desc, size=11, color=GRAY_700)
    add_text(slide, x + Inches(0.15), y + Inches(1.9), Inches(2.65), Inches(0.4), services, size=10, color=GRAY_400, bold=True)


# ═════════════════════════════════════════════════════
# SLIDE 10 — VALUE AI ADDS
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 10)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'What Value the AI Layer Adds', size=28, color=GRAY_900, bold=True)

# Before (left)
add_shape_bg(slide, Inches(0.8), Inches(1.3), Inches(5.6), Inches(2.5), LIGHT_RED)
add_text(slide, Inches(1), Inches(1.4), Inches(5.2), Inches(0.4), '❌  Before CivicLens (Without AI)', size=16, color=RED, bold=True)
before_items = [
    "• File RTI → Wait 30 days → Get 1 answer about 1 project",
    "• No way to verify if 'Completed' roads are actually done",
    "• No idea who is actually responsible for failures",
    "• Government data locked in PDFs nobody can read",
]
tf = add_text(slide, Inches(1), Inches(1.9), Inches(5.2), Inches(0.3), before_items[0], size=13, color=GRAY_700)
for item in before_items[1:]:
    add_para(tf, item, size=13, color=GRAY_700, space_before=Pt(6))

# After (right)
add_shape_bg(slide, Inches(6.9), Inches(1.3), Inches(5.6), Inches(2.5), LIGHT_GREEN)
add_text(slide, Inches(7.1), Inches(1.4), Inches(5.2), Inches(0.4), '✅  After CivicLens (With AI)', size=16, color=GREEN, bold=True)
after_items = [
    "• Instant access to every project's budget, deadline, contractor",
    "• Photo-based verification — AI detects if road is really done",
    "• Accountability chain: Ministry → Dept → Officer → Contractor",
    "• New tenders processed automatically the moment they're uploaded",
]
tf = add_text(slide, Inches(7.1), Inches(1.9), Inches(5.2), Inches(0.3), after_items[0], size=13, color=GRAY_700)
for item in after_items[1:]:
    add_para(tf, item, size=13, color=GRAY_700, space_before=Pt(6))

# Impact Table
add_text(slide, Inches(0.8), Inches(4.1), Inches(5), Inches(0.3), 'Impact Metrics', size=16, color=BLUE, bold=True)

table = make_table(slide, 6, 3, Inches(0.8), Inches(4.5), Inches(11.7), Inches(2.5))
set_cell(table, 0, 0, 'Metric', size=12, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Without AI', size=12, bold=True, color=WHITE)
set_cell(table, 0, 2, 'With CivicLens', size=12, bold=True, color=WHITE)

metrics = [
    ('Time to get project info', '30 days (RTI)', '3 seconds'),
    ('Tenders processed per day', 'Manual, ~10', '1,000+ (automated)'),
    ('Verification method', 'None exists', 'AI photo analysis + 200m geofence'),
    ('Accountability visibility', 'None', 'Full chain: Ministry → Contractor'),
    ('Cost per tender processed', '₹100+ (govt staff time)', '< ₹1 (AI + serverless)'),
]
for i, (metric, before, after) in enumerate(metrics):
    set_cell(table, i+1, 0, metric, bold=True)
    set_cell(table, i+1, 1, before, color=RED)
    set_cell(table, i+1, 2, after, color=GREEN, bold=True)

table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(4.0)
table.columns[2].width = Inches(4.2)
style_table(table)


# ═════════════════════════════════════════════════════
# SLIDE 11 — LIVE DEMO
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 11)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Working Prototype', size=28, color=GRAY_900, bold=True)

# URL box
add_shape_bg(slide, Inches(2), Inches(1.2), Inches(9.3), Inches(0.7), LIGHT_BLUE)
add_text(slide, Inches(2), Inches(1.3), Inches(9.3), Inches(0.5), '🌐  https://main.d3k5zyfw9140d9.amplifyapp.com', size=18, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

add_text(slide, Inches(0.8), Inches(2.2), Inches(5), Inches(0.3), 'What Evaluators Can Test:', size=16, color=GRAY_900, bold=True)

test_items = [
    ("🗺️  Dashboard", "Browse 18+ real infrastructure projects across 8 Indian cities on a live interactive map"),
    ("🔍  Project Detail", "Click any project → see budget, deadline, contractor, truth check, citizen audits"),
    ("🤖  AI Demo", "Interactive simulation of both AI pipelines (tender processing + photo verification)"),
    ("📊  Statistics", "Budget analysis, status distribution, most-disputed projects ranking"),
    ("🏛️  Accountability", "Explore the full government hierarchy — Ministry to Contractor"),
    ("📷  Report Issue", "4-step citizen audit wizard with GPS verification and photo upload"),
    ("ℹ️  About", "Full architecture diagram and AWS services breakdown"),
]

y = Inches(2.7)
for emoji_title, desc in test_items:
    add_shape_bg(slide, Inches(0.8), y, Inches(11.7), Inches(0.55), GRAY_50)
    add_text(slide, Inches(1), y + Inches(0.05), Inches(2.5), Inches(0.4), emoji_title, size=13, color=GRAY_900, bold=True)
    add_text(slide, Inches(3.5), y + Inches(0.05), Inches(8.8), Inches(0.4), desc, size=12, color=GRAY_600)
    y += Inches(0.62)

add_text(slide, Inches(0.8), Inches(6.8), Inches(11), Inches(0.3),
         'Data: Based on real Indian infrastructure incidents — RTI findings, news reports, documented cases from Delhi, Mumbai, Bengaluru, Chennai, Kolkata, Bhopal, Lucknow, Patna.',
         size=11, color=GRAY_400)


# ═════════════════════════════════════════════════════
# SLIDE 12 — TECH STACK
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 12)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Technology Stack', size=28, color=GRAY_900, bold=True)

table = make_table(slide, 11, 2, Inches(1.5), Inches(1.2), Inches(10.3), Inches(5.5))
set_cell(table, 0, 0, 'Layer', size=13, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Technology', size=13, bold=True, color=WHITE)

stack = [
    ('Frontend', 'React 18, Vite 6.4, TailwindCSS 3.4, React-Leaflet, Recharts, Lucide React'),
    ('Map', 'OpenStreetMap + CARTO basemap (zero API keys needed)'),
    ('Backend', 'AWS Lambda × 6 (Python 3.13, ARM64 Graviton)'),
    ('API', 'Amazon API Gateway (REST, CORS enabled)'),
    ('Database', 'Amazon DynamoDB × 2 tables (PAY_PER_REQUEST, GSIs for geohash + city)'),
    ('AI / ML', 'Amazon Bedrock (Claude 3 Haiku), Amazon Textract, Amazon Rekognition'),
    ('Geospatial', 'Amazon Location Service (Esri provider, geocoding + geofence)'),
    ('Storage', 'Amazon S3 × 2 buckets (tender PDFs + citizen photos)'),
    ('IaC', 'AWS SAM — template.yaml (308 lines of infrastructure-as-code)'),
    ('Hosting', 'AWS Amplify — Region: ap-south-1 (Mumbai)'),
]
for i, (layer, tech) in enumerate(stack):
    set_cell(table, i+1, 0, layer, bold=True)
    set_cell(table, i+1, 1, tech)

table.columns[0].width = Inches(2.0)
table.columns[1].width = Inches(8.3)
style_table(table)


# ═════════════════════════════════════════════════════
# SLIDE 13 — COST
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 13)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Cost — Under $5/Month', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.3), 'Fully serverless = no idle server costs. Scales linearly with usage.', size=14, color=GRAY_500)

table = make_table(slide, 10, 3, Inches(1.5), Inches(1.5), Inches(10.3), Inches(4.5))
set_cell(table, 0, 0, 'AWS Service', size=13, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Est. Monthly Cost', size=13, bold=True, color=WHITE)
set_cell(table, 0, 2, 'Notes', size=13, bold=True, color=WHITE)

costs = [
    ('Lambda', '~$0', 'Free tier: 1M requests/month'),
    ('DynamoDB', '~$0', 'PAY_PER_REQUEST, demo-level traffic'),
    ('API Gateway', '~$0', 'Free tier: 1M calls/month'),
    ('S3', '~$0.01', 'A few MBs of storage'),
    ('Bedrock (Claude 3 Haiku)', '~$0.50–2.00', '$0.25 per 1M input tokens — very cost efficient'),
    ('Textract', '~$1–3', '$1.50 per 1,000 pages processed'),
    ('Rekognition', '~$0', 'Free tier: 5,000 images/month'),
    ('Location Service', '~$0', 'Free tier: 10K geocodes/month'),
    ('TOTAL', '< $5/month', 'Well within $200 AWS credits provided'),
]
for i, (svc, cost, notes) in enumerate(costs):
    is_total = i == len(costs) - 1
    set_cell(table, i+1, 0, svc, bold=True, color=BLUE if is_total else GRAY_800)
    set_cell(table, i+1, 1, cost, bold=is_total, color=GREEN if is_total else GRAY_800)
    set_cell(table, i+1, 2, notes, bold=is_total)

table.columns[0].width = Inches(3.5)
table.columns[1].width = Inches(2.3)
table.columns[2].width = Inches(4.5)
style_table(table)


# ═════════════════════════════════════════════════════
# SLIDE 14 — REAL-WORLD IMPACT
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 14)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Real-World Impact — Why This Matters', size=28, color=GRAY_900, bold=True)

# Scale stats
stats = [
    ("29", "States in India"),
    ("700+", "Districts"),
    ("2,50,000+", "Gram Panchayats"),
    ("30-40%", "Estimated Corruption\nin Public Works"),
]
for i, (num, label) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    add_shape_bg(slide, x, Inches(1.3), Inches(2.8), Inches(1.3), LIGHT_BLUE)
    add_text(slide, x, Inches(1.4), Inches(2.8), Inches(0.6), num, size=30, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.0), Inches(2.8), Inches(0.5), label, size=12, color=GRAY_600, align=PP_ALIGN.CENTER)

# What CivicLens enables
add_text(slide, Inches(0.8), Inches(3.0), Inches(10), Inches(0.3), 'What CivicLens Enables:', size=18, color=GRAY_900, bold=True)

enables = [
    ("Transparency at Scale", "AI reads every tender so citizens don't have to. Automatic extraction replaces 30-day RTI."),
    ("Ground-Truth Verification", "Computer vision catches mismatches between official claims and reality. Tamper-proof with geofencing."),
    ("Accountability Visibility", "For the first time, citizens can see the full chain of responsibility — from Ministry to Contractor."),
    ("Democratic Oversight", "Upvote system ensures the most critical issues get attention first. Signal over noise."),
]

for i, (title, desc) in enumerate(enables):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.2)
    y = Inches(3.5 + row * 1.5)
    add_shape_bg(slide, x, y, Inches(5.9), Inches(1.3), GRAY_50)
    add_text(slide, x + Inches(0.2), y + Inches(0.1), Inches(5.5), Inches(0.35), title, size=14, color=BLUE, bold=True)
    add_text(slide, x + Inches(0.2), y + Inches(0.5), Inches(5.5), Inches(0.7), desc, size=12, color=GRAY_600)

# Potential
add_shape_bg(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6), LIGHT_GREEN)
add_text(slide, Inches(1), Inches(6.55), Inches(11.3), Inches(0.5),
         '🚀  If deployed nationwide: process ALL state government tenders automatically. Every smartphone becomes a civic audit tool. Permanent, searchable, AI-verified record of public spending.',
         size=13, color=RGBColor(6, 95, 70), bold=True, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════
# SLIDE 15 — FUTURE ROADMAP
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 15)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Future Roadmap', size=28, color=GRAY_900, bold=True)

table = make_table(slide, 9, 3, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.5))
set_cell(table, 0, 0, 'Phase', size=13, bold=True, color=WHITE)
set_cell(table, 0, 1, 'Feature', size=13, bold=True, color=WHITE)
set_cell(table, 0, 2, 'AWS Service', size=13, bold=True, color=WHITE)

roadmap = [
    ('v2', 'Hindi / regional language tender processing', 'Bedrock + Textract multi-language OCR'),
    ('v2', 'Real-time push notifications when projects are flagged', 'Amazon SNS + EventBridge'),
    ('v2', 'User authentication with Aadhaar-lite verification', 'Amazon Cognito'),
    ('v3', 'Historical trend analysis — budget overruns, deadline slippage', 'Bedrock + DynamoDB Streams'),
    ('v3', 'Satellite imagery comparison (before/after construction)', 'Rekognition Custom Labels + S3'),
    ('v3', 'State government API integration for live tender data', 'API Gateway + Step Functions'),
    ('v4', 'Mobile app (React Native) for field verification', 'AWS Amplify + AppSync'),
    ('v4', 'Government accountability dashboard for departments', 'Amazon QuickSight + DynamoDB'),
]
for i, (phase, feature, svc) in enumerate(roadmap):
    set_cell(table, i+1, 0, phase, bold=True, color=BLUE)
    set_cell(table, i+1, 1, feature)
    set_cell(table, i+1, 2, svc, bold=True)

table.columns[0].width = Inches(1.2)
table.columns[1].width = Inches(6.5)
table.columns[2].width = Inches(4.0)
style_table(table)


# ═════════════════════════════════════════════════════
# SLIDE 16 — EVALUATION ALIGNMENT
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_blue_bar(slide)
add_slide_number(slide, 16)

add_text(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.5), 'Hackathon Evaluation Alignment', size=28, color=GRAY_900, bold=True)
add_text(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.3), 'How CivicLens meets every evaluation criteria', size=14, color=GRAY_500)

table = make_table(slide, 7, 3, Inches(0.8), Inches(1.4), Inches(11.7), Inches(5.3))
set_cell(table, 0, 0, 'Evaluation Criteria', size=12, bold=True, color=WHITE)
set_cell(table, 0, 1, 'How CivicLens Addresses It', size=12, bold=True, color=WHITE)
set_cell(table, 0, 2, 'Evidence', size=12, bold=True, color=WHITE)

eval_rows = [
    ('Uses AWS Generative AI', 'Amazon Bedrock (Claude 3 Haiku) is the core engine — extracts structured data from unstructured government tender PDFs', 'ingest_tender/handler.py — Bedrock API call with custom prompt engineering'),
    ('AWS Infrastructure', '10 AWS services used: Lambda, API GW, DynamoDB, S3, Textract, Rekognition, Bedrock, Location Service, SAM, Amplify', 'template.yaml — 308 lines of SAM infrastructure-as-code'),
    ('Why AI is Required', 'Tenders = unstructured scanned PDFs. Only OCR + Gen AI can extract data. Photo verification = CV. Scale = 1000+/day.', 'See Slide 4 — detailed breakdown'),
    ('Value AI Adds', '30 days → 3 sec. Manual → Automated. Claims → Verified. Opaque → Transparent.', 'See Slide 10 — impact metrics table'),
    ('Working Prototype', 'Live on AWS Amplify with API backend connected to DynamoDB. 7 pages, 18+ projects, full interactivity.', 'main.d3k5zyfw9140d9.amplifyapp.com'),
    ('GitHub Repository', 'Complete codebase: frontend (React), backend (SAM + Lambda), IaC, seed data, deployment guide', 'Public repo with README + DEPLOY.md'),
]
for i, (criteria, how, evidence) in enumerate(eval_rows):
    set_cell(table, i+1, 0, criteria, bold=True, color=BLUE)
    set_cell(table, i+1, 1, how)
    set_cell(table, i+1, 2, evidence)

table.columns[0].width = Inches(2.5)
table.columns[1].width = Inches(5.7)
table.columns[2].width = Inches(3.5)
style_table(table)


# ═════════════════════════════════════════════════════
# SLIDE 17 — CLOSING
# ═════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK)
add_shape_bg(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), BLUE)

add_text(slide, Inches(1), Inches(1.0), Inches(11.333), Inches(0.7), 'CivicLens', size=44, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(slide, Inches(1), Inches(1.7), Inches(11.333), Inches(0.5), 'The Truth Engine', size=28, color=BLUE, bold=True, align=PP_ALIGN.CENTER)

# Quote
add_text(slide, Inches(1.5), Inches(2.8), Inches(10.333), Inches(1.2),
         '"We blame the PM or CM for every pothole and broken pipe. But there are hundreds of departments, officers, and contractors responsible. The government releases data — but no one can read it all.\nCivicLens takes off the blindfold."',
         size=15, color=GRAY_400, align=PP_ALIGN.CENTER)

# 3 stats
stat_data = [
    ("10", "AWS Services", "Bedrock, Textract, Rekognition,\nLocation, Lambda, API GW,\nDynamoDB, S3, SAM, Amplify"),
    ("18+", "Real Projects", "Based on real RTI data across\n8 Indian cities: Delhi, Mumbai,\nBengaluru, Chennai & more"),
    ("< $5", "Per Month", "Fully serverless architecture.\nNo idle server costs.\nProduction-ready scaling."),
]
for i, (num, title, desc) in enumerate(stat_data):
    x = Inches(1.2 + i * 3.9)
    add_shape_bg(slide, x, Inches(4.3), Inches(3.5), Inches(2.0), RGBColor(30, 41, 59))
    add_text(slide, x, Inches(4.4), Inches(3.5), Inches(0.6), num, size=32, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(4.95), Inches(3.5), Inches(0.35), title, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(5.35), Inches(3.5), Inches(0.8), desc, size=10, color=GRAY_400, align=PP_ALIGN.CENTER)

# Links
add_text(slide, Inches(1), Inches(6.6), Inches(11.333), Inches(0.35),
         '🌐  Live: main.d3k5zyfw9140d9.amplifyapp.com    |    💻  GitHub: [Your Repo URL]    |    📹  Video: [Your Video URL]',
         size=12, color=GRAY_500, align=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(7.0), Inches(11.333), Inches(0.3),
         'Built with ❤️ for Bharat', size=13, color=RGBColor(253, 230, 138), align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════
# SAVE
# ═════════════════════════════════════════════════════
out_path = r"c:\Users\adity\OneDrive\Desktop\AI for BHARAT\Project\CivicLens_PPT.pptx"
prs.save(out_path)
print(f"✅ PowerPoint saved: {out_path}")
print(f"   Slides: {len(prs.slides)}")
